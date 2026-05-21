"""文档解析器 - 通用工具类,支持多种文件格式解析"""
from typing import List
from pathlib import Path
from loguru import logger
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    Docx2txtLoader,
    CSVLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter


class SimpleDocParser:
    """基于LangChain的文档解析器"""
    
    # 文本分割器配置 - 可配置化
    CHUNK_SIZE = 500
    CHUNK_OVERLAP = 50
    
    # 支持的文件格式映射
    SUPPORTED_FORMATS = {
        '.txt': '纯文本',
        '.md': 'Markdown',
        '.pdf': 'PDF文档',
        '.docx': 'Word文档',
        '.html': 'HTML页面',
        '.htm': 'HTML页面',
        '.csv': 'CSV表格',
        '.json': 'JSON数据',
        '.xlsx': 'Excel表格',
        '.xls': 'Excel表格(旧版)',
        '.pptx': 'PowerPoint演示文稿',
        '.ppt': 'PowerPoint演示文稿(旧版)',
        '.rtf': '富文本格式',
        '.xml': 'XML数据',
    }
    
    @classmethod
    def _get_text_splitter(cls, chunk_size: int = None, chunk_overlap: int = None) -> RecursiveCharacterTextSplitter:
        """
        获取文本分割器
        
        Args:
            chunk_size: 分块大小，默认使用类配置
            chunk_overlap: 分块重叠大小，默认使用类配置
            
        Returns:
            RecursiveCharacterTextSplitter实例
        """
        size = chunk_size or cls.CHUNK_SIZE
        overlap = chunk_overlap or cls.CHUNK_OVERLAP
        
        return RecursiveCharacterTextSplitter(
            chunk_size=size,
            chunk_overlap=overlap,
            length_function=len,
            separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
        )
    
    @staticmethod
    def parse_text(file_path: str) -> List[str]:
        """
        解析文本文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            loader = TextLoader(file_path, encoding='utf-8')
            documents = loader.load()
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents(documents)
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析文本文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析文本文件失败: {e}")
            return []
    
    @staticmethod
    def parse_markdown(file_path: str) -> List[str]:
        """
        解析Markdown文件（使用简单文本加载器）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            # 使用TextLoader读取Markdown文件内容
            loader = TextLoader(file_path, encoding='utf-8')
            documents = loader.load()
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents(documents)
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析Markdown文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析Markdown文件失败: {e}")
            return []
    
    @staticmethod
    def parse_pdf(file_path: str) -> List[str]:
        """
        解析PDF文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents(documents)
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析PDF文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析PDF文件失败: {e}")
            return []
    
    @staticmethod
    def parse_word(file_path: str) -> List[str]:
        """
        解析Word文档 (.docx)
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents(documents)
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析Word文档 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析Word文档失败: {e}")
            return []
    
    @staticmethod
    def parse_html(file_path: str) -> List[str]:
        """
        解析HTML文件（使用BeautifulSoup提取文本）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用BeautifulSoup提取纯文本
            soup = BeautifulSoup(html_content, 'html.parser')
            text = soup.get_text(separator='\n', strip=True)
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=text)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析HTML文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析HTML文件失败: {e}")
            return []
    
    @staticmethod
    def parse_csv(file_path: str) -> List[str]:
        """
        解析CSV文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            loader = CSVLoader(file_path, encoding='utf-8')
            documents = loader.load()
            
            # CSV每行是一个document，直接拼接
            chunk_texts = [
                doc.page_content.strip() 
                for doc in documents 
                if doc.page_content.strip()
            ]
            
            logger.info(f"解析CSV文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析CSV文件失败: {e}")
            return []
    
    @staticmethod
    def parse_json(file_path: str) -> List[str]:
        """
        解析JSON文件（简单实现）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # 将JSON转换为字符串
            if isinstance(data, (dict, list)):
                text_content = json.dumps(data, ensure_ascii=False, indent=2)
            else:
                text_content = str(data)
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=text_content)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析JSON文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析JSON文件失败: {e}")
            return []
    
    @staticmethod
    def parse_excel(file_path: str) -> List[str]:
        """
        解析Excel文件 (.xlsx, .xls)
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            import pandas as pd
            
            # 读取所有sheet
            if file_path.endswith('.xlsx'):
                xl = pd.ExcelFile(file_path)
                sheets = xl.sheet_names
            else:
                # .xls格式
                xl = pd.ExcelFile(file_path)
                sheets = xl.sheet_names
            
            all_text = []
            for sheet_name in sheets:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_text = f"Sheet: {sheet_name}\n" + df.to_string(index=False)
                all_text.append(sheet_text)
            
            full_content = "\n\n".join(all_text)
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=full_content)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析Excel文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析Excel文件失败: {e}")
            return []
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[str]:
        """
        解析PowerPoint文件 (.pptx)
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_content = f"Slide {slide_idx}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content += shape.text + "\n"
                slides_text.append(slide_content.strip())
            
            full_content = "\n\n".join(slides_text)
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=full_content)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析PowerPoint文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析PowerPoint文件失败: {e}")
            return []
    
    @staticmethod
    def parse_rtf(file_path: str) -> List[str]:
        """
        解析RTF文件 (富文本格式)
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            # 使用striprtf库解析RTF
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=text)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析RTF文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析RTF文件失败: {e}")
            return []
    
    @staticmethod
    def parse_xml(file_path: str) -> List[str]:
        """
        解析XML文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文本块列表
        """
        try:
            import xml.etree.ElementTree as ET
            
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # 提取XML中的所有文本内容
            text_content = ET.tostring(root, encoding='unicode', method='text')
            
            # 创建临时文档对象
            from langchain_core.documents import Document
            doc = Document(page_content=text_content)
            
            text_splitter = SimpleDocParser._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            
            logger.info(f"解析XML文件 {file_path}，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
            
        except Exception as e:
            logger.error(f"解析XML文件失败: {e}")
            return []
    
    @classmethod
    def parse_content(cls, content: str) -> List[str]:
        """
        直接解析文本内容
        
        Args:
            content: 文本内容
            
        Returns:
            文本块列表
        """
        try:
            from langchain_core.documents import Document
            doc = Document(page_content=content)
            
            text_splitter = cls._get_text_splitter()
            chunks = text_splitter.split_documents([doc])
            
            chunk_texts = [chunk.page_content.strip() for chunk in chunks if chunk.page_content.strip()]
            logger.info(f"解析文本内容，得到 {len(chunk_texts)} 个文本块")
            return chunk_texts
        except Exception as e:
            logger.error(f"解析文本内容失败: {e}")
            return []

    @classmethod
    def parse_file(cls, file_path: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """
        根据文件类型解析文件
        
        Args:
            file_path: 文件路径
            chunk_size: 分块大小，可选
            chunk_overlap: 分块重叠大小，可选
            
        Returns:
            文本块列表
        """
        suffix = Path(file_path).suffix.lower()
        
        # 根据文件类型调整分块策略
        if chunk_size is None:
            # PDF和Word使用较小的分块
            if suffix in ['.pdf', '.docx']:
                chunk_size = 400
            # Excel和PPT使用中等分块
            elif suffix in ['.xlsx', '.pptx']:
                chunk_size = 600
            else:
                chunk_size = cls.CHUNK_SIZE
        
        if chunk_overlap is None:
            chunk_overlap = cls.CHUNK_OVERLAP
        
        # 临时设置分块参数
        old_chunk_size = cls.CHUNK_SIZE
        old_chunk_overlap = cls.CHUNK_OVERLAP
        cls.CHUNK_SIZE = chunk_size
        cls.CHUNK_OVERLAP = chunk_overlap
        
        try:
            # 文本文件
            if suffix in ['.txt', '.text']:
                return cls.parse_text(file_path)
            # Markdown
            elif suffix == '.md':
                return cls.parse_markdown(file_path)
            # PDF
            elif suffix == '.pdf':
                return cls.parse_pdf(file_path)
            # Word文档
            elif suffix in ['.docx', '.doc']:
                if suffix == '.doc':
                    logger.warning(".doc格式不支持，请使用.docx格式")
                    return []
                return cls.parse_word(file_path)
            # HTML
            elif suffix in ['.html', '.htm']:
                return cls.parse_html(file_path)
            # CSV
            elif suffix == '.csv':
                return cls.parse_csv(file_path)
            # JSON
            elif suffix == '.json':
                return cls.parse_json(file_path)
            # Excel
            elif suffix in ['.xlsx', '.xls']:
                return cls.parse_excel(file_path)
            # PowerPoint
            elif suffix in ['.pptx', '.ppt']:
                if suffix == '.ppt':
                    logger.warning(".ppt格式不支持，请使用.pptx格式")
                    return []
                return cls.parse_pptx(file_path)
            # RTF
            elif suffix == '.rtf':
                return cls.parse_rtf(file_path)
            # XML
            elif suffix == '.xml':
                return cls.parse_xml(file_path)
            else:
                logger.warning(f"不支持的文件类型: {suffix}")
                return []
        finally:
            # 恢复原始配置
            cls.CHUNK_SIZE = old_chunk_size
            cls.CHUNK_OVERLAP = old_chunk_overlap
